# -*- conding: utf-8 -*-

import Reader
from pptx import Presentation
from pptx.shapes.picture import Picture
from os.path import basename
import imghdr
from datetime import datetime
Img = ['jpg', 'JPG', 'jpeg' , 'JPEG' , 'png', 'PNG', 'gif',  'GIF', 'bmp', 'BMP', 'tif', 'TIF', 'tiff', 'TIFF']
index = 1
class PPtReader:

    def __init__(self, src: str, office_list: [], imgdir: str) -> None:
        self.src = src
        self.office_list = office_list
        self.ImageDir = imgdir

    def getText(self): #https://cloud.tencent.com/developer/article/1708628
        res = []
        res_img = []
        img_name = []
        c = 0
        for file in self.office_list:
            if file.split('.')[-1] == 'pptx':
                presentation = Presentation(file)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            res.append(text_frame.text+"\n")
                        elif shape.has_table:
                            for row in shape.table.rows: #读每行
                                for cell in row.cells: #读一行的所有单元格
                                    res.append(cell.text)
                        elif isinstance(shape, Picture):
                            #shape.image.blob #图像二进制字节流
                            imagetype = shape.image.content_type
                            imtype = imagetype[imagetype.find('/') +1 : ] #后缀名\
                            # imtype = '.png'
                            
                            if imtype not in Img:
                                continue
                            file = file.split('/')[-1]
                            suffix = datetime.strftime(datetime.now(),'%Y%m%d-%H%M%S')+ '.png'
                            suffix = str(c)+'.png'
                            c += 1
                            img_name.append(file+suffix)
                            res_img.append(shape.image.blob)
        with open('fileDIR/ppt_text.txt', 'w', encoding='utf-8') as f:
            for i in res:
                f.write(i)
        c = 0
        for i in img_name:
            with open(self.ImageDir + i, "wb") as f:
            # with open('../IMAGE/'+i, "wb") as f:
                f.write(res_img[c])
            c += 1


if __name__ == '__main__':
    R = Reader.Reader(u'dDIR/')
    R.file_reader()
    pr = PPtReader('dDIR/', R.office_list)
    pr.getText()
    # pr.getPicture()
