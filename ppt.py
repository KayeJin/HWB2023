# -*- conding: utf-8 -*-

import Reader
from pptx import Presentation
from pptx.shapes.picture import Picture

index = 1
class PPtReader:

    def __init__(self, src: str, office_list: []) -> None:
        self.src = src
        self.office_list = office_list

    def getText(self): #https://cloud.tencent.com/developer/article/1708628
        print(self.office_list)
        res = []
        for file in self.office_list:
            if file.split('.')[1] == 'pptx':
                presentation = Presentation(self.src+"/"+file)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            res.append(text_frame.text+"\n")
                            # print(text_frame.text)
        
        with open('../ppt_text', 'w', encoding='utf-8') as f:
            for i in res:
                f.write(i)

    def getTable(self):
        res = []
        for file in self.office_list:
            if file.split('.')[1] == 'pptx':
                presentation = Presentation(self.src+"/"+file)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_table:
                            for row in shape.table.rows: #读每行
                                for cell in row.cells: #读一行的所有单元格
                                    res.append(cell.text)

    def getPicture(self):
        global index
        for file in self.office_list:
            if file.split('.')[1] == 'pptx':
                presentation = Presentation(self.src+"/"+file)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if isinstance(shape, Picture):
                            #shape.image.blob #图像二进制字节流
                            with open(f'../ppt_{index}.jpg', 'wb') as f:
                                f.write(shape.image.blob)
                                index += 1




if __name__ == '__main__':
    R = Reader.Reader(u'../赛题材料/office')
    R.file_reader()
    pr = PPtReader("../赛题材料/office", R.office_list)
    pr.getText()
    pr.getPicture()
